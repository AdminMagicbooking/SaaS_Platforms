# -*- coding: utf-8 -*-
from pptx import Presentation

# (prefix_to_match, new_full_text)
COMMON = [
    ("ECOGRID RESORT™  ×  ETHICAL SOLAR", "ECOGRID RESORT™"),
]
FR = COMMON + [
    ("Préparé pour ETHICAL SOLAR", "Document confidentiel  ·  par l’équipe COREPROMA  ·  2026"),
    ("Une JV de co-développement", "Co-développement & co-investissement — un modèle ouvert"),
    ("Ethical Solar apporte", "L’investisseur — trois profils ouverts"),
    ("Co-investit en fonds propres", "Industriel / EPC : capital, installation & marge EPC"),
    ("Pilote l'EPC", "Infrastructure : ombrières PV abritant onduleurs & batteries"),
    ("Perçoit la marge EPC", "Financier (fonds / VC) : apport en capital, rendement"),
    ("Mix de financement pilote", "Selon le profil : co-investissement en fonds propres par SPV d’actifs (un par site ou grappe) et/ou apport industriel (EPC, ombrières). Cœur logiciel COREPROMA partagé. Mix indicatif du pilote (5,7 M€) : 1,4 M€ fonds propres · 3,4 M€ dette verte senior · 0,6 M€ aides · 0,3 M€ leasing."),
    ("Ce que nous demandons à Ethical Solar", "Ce que nous attendons d’un partenaire investisseur"),
]
EN = COMMON + [
    ("Prepared for ETHICAL SOLAR", "Confidential document  ·  by the COREPROMA team  ·  2026"),
    ("A co-development & co-investment JV", "Co-development & co-investment — an open model"),
    ("Ethical Solar brings", "The investor — three open profiles"),
    ("Co-invest equity in each asset SPV", "Industrial / EPC: capital, build & EPC margin"),
    ("Lead EPC / build", "Infrastructure: PV carports housing inverters & batteries"),
    ("Earn EPC margin", "Financial (fund / VC): capital only, financial return"),
    ("Pilot funding mix", "Depending on profile: equity co-investment via asset SPVs (one per site or cluster) and/or industrial contribution (EPC, carports). Shared COREPROMA software core. Indicative pilot mix (5.7 M€): 1.4 M€ equity · 3.4 M€ green senior debt · 0.6 M€ grants · 0.3 M€ leasing."),
    ("The ask of Ethical Solar", "What we expect from an investor partner"),
]

def apply(path, rules):
    prs = Presentation(path)
    hits = {}
    for s in prs.slides:
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                txt = "".join(r.text for r in para.runs).strip()
                if not txt:
                    continue
                for pref, new in rules:
                    if txt.startswith(pref):
                        if para.runs:
                            para.runs[0].text = new
                            for r in para.runs[1:]:
                                r.text = ""
                        hits[pref] = hits.get(pref, 0) + 1
                        break
    prs.save(path)
    return hits

for path, rules in [("ECOGRID_RESORT_x_Ethical_Solar_FR.pptx", FR), ("ECOGRID_RESORT_x_Ethical_Solar_EN.pptx", EN)]:
    h = apply(path, rules)
    print(path)
    for pref, n in h.items():
        print("   ", n, "x", pref[:45])
    missing = [p for p, _ in rules if p not in h]
    if missing:
        print("   !! NOT MATCHED:", missing)
