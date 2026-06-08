from pptx import Presentation
from pptx.util import Pt
import copy, re

def update_slide2_en(prs):
    slide = prs.slides[1]  # slide index 1 = slide 2
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                t = run.text

                # 1. Expand energy production bullet to include wind + V2G
                if 'On-site solar PV: car-park shades, roofs, pergolas' in t:
                    run.text = t.replace(
                        'On-site solar PV: car-park shades, roofs, pergolas',
                        'On-site solar PV: carports, roofs, pergolas'
                    )

                # 2. Update "Store & manage" section to include V2G framing
                if '4–8 MWh LFP batteries: resilience + arbitrage' in t:
                    run.text = t.replace(
                        '4–8 MWh LFP batteries: resilience + arbitrage',
                        '4–8 MWh LFP batteries (V2G-ready): resilience, arbitrage & local network feed'
                    )

                if 'AI EMS forecasts and dispatches 24/7' in t:
                    run.text = t.replace(
                        'AI EMS forecasts and dispatches 24/7',
                        'AI EMS forecasts and dispatches 24/7 — including wind & V2G sources'
                    )

                # 3. Add wind turbine reference to Produce section
                if 'Turns idle surfaces into productive assets' in t:
                    run.text = t.replace(
                        'Turns idle surfaces into productive assets',
                        'Turns idle surfaces into productive assets. Wind turbine (optional complement where available)'
                    )

                # 4. Update EV chargers to show V2G
                if '30–50 EV charge points, smart charging' in t:
                    run.text = t.replace(
                        '30–50 EV charge points, smart charging',
                        '30–50 EV charge points — smart charging + V2G feedback to local network'
                    )

                # 5. Expand guest app bullet
                if 'Guest app: EcoScore, EV booking, rewards' in t:
                    run.text = t.replace(
                        'Guest app: EcoScore, EV booking, rewards',
                        'Guest app: EcoScore, real-time savings & CO₂ tracker, EV booking, V2G rewards — personalised stay report at checkout'
                    )

                # 6. Update "Engage & report" label if present
                if 'Engage & report' in t:
                    run.text = t.replace('Engage & report', 'Engage & measure')

    return prs


def update_slide2_fr(prs):
    slide = prs.slides[1]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                t = run.text

                if 'Solaire sur site : ombrières, toitures, pergolas' in t:
                    run.text = t.replace(
                        'Solaire sur site : ombrières, toitures, pergolas',
                        'Solaire sur site : ombrières, toitures, pergolas. Éolienne (option, si disponible sur site)'
                    )

                if 'Batteries LFP' in t and 'résilience' in t:
                    run.text = t.replace(
                        'Batteries LFP (résilience + arbitrage)',
                        'Batteries LFP V2G-prêtes (résilience, arbitrage & injection réseau local)'
                    ).replace(
                        'Batteries LFP',
                        'Batteries LFP V2G-prêtes (résilience, arbitrage & injection réseau local) ;'
                    )

                if '30–50 bornes VE' in t and 'recharge intelligente' in t:
                    run.text = t.replace(
                        '30–50 bornes VE, recharge intelligente',
                        '30–50 bornes VE — recharge intelligente + retour V2G réseau local'
                    )

                if 'App client : EcoScore' in t:
                    run.text = t.replace(
                        'App client : EcoScore, réservation VE, récompenses',
                        'App vacancier : EcoScore, suivi économies & CO₂ en temps réel, réservation VE, récompenses V2G — bilan personnalisé au départ'
                    )

    return prs


# Process EN deck
prs_en = Presentation('/sessions/kind-exciting-edison/mnt/SaaS_Projects/SmartGrid/Commercial/ECOGRID_RESORT_Investors_EN.pptx')
prs_en = update_slide2_en(prs_en)
prs_en.save('/sessions/kind-exciting-edison/mnt/SaaS_Projects/SmartGrid/Commercial/ECOGRID_RESORT_Investors_EN.pptx')
print("EN deck updated")

# Process FR deck
prs_fr = Presentation('/sessions/kind-exciting-edison/mnt/SaaS_Projects/SmartGrid/Commercial/ECOGRID_RESORT_Investisseurs_FR.pptx')
prs_fr = update_slide2_fr(prs_fr)
prs_fr.save('/sessions/kind-exciting-edison/mnt/SaaS_Projects/SmartGrid/Commercial/ECOGRID_RESORT_Investisseurs_FR.pptx')
print("FR deck updated")
